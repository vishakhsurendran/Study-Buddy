import os
import logging
from pathlib import Path
from typing import List, Dict, Any

import fitz  # pymupdf
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

class ResourceIntake:
    @staticmethod
    def simple_chunker(text: str, max_words: int = 200, overlap: int = 40) -> List[str]:
        words = text.split()
        if len(words) <= max_words:
            return [text]
        chunks = []
        start = 0
        while start < len(words):
            end = min(start + max_words, len(words))
            chunks.append(" ".join(words[start:end]))
            if end == len(words):
                break
            start = end - overlap
        return chunks

    @staticmethod
    def extract_docx(path: str, chunk_words: int = 200, overlap: int = 40) -> List[Dict[str, Any]]:
        results = []
        doc = DocxDocument(path)
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            base_meta = {"source": os.path.basename(path), "type": "docx", "para_idx": i}
            for chunk_idx, sub in enumerate(ResourceIntake.simple_chunker(text, chunk_words, overlap)):
                meta = dict(base_meta)
                meta.update({"chunk_idx": chunk_idx, "excerpt": sub[:200]})
                results.append({"text": sub, "meta": meta})
        return results

    @staticmethod
    def extract_pptx(path: str, chunk_words: int = 200, overlap: int = 40) -> List[Dict[str, Any]]:
        results = []
        prs = Presentation(path)
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue
                    base_meta = {
                        "source": os.path.basename(path),
                        "type": "pptx",
                        "slide_idx": slide_idx,
                        "shape_idx": shape_idx,
                    }
                    for chunk_idx, sub in enumerate(ResourceIntake.simple_chunker(text, chunk_words, overlap)):
                        meta = dict(base_meta)
                        meta.update({"chunk_idx": chunk_idx, "excerpt": sub[:200]})
                        results.append({"text": sub, "meta": meta})
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        base_meta = {"source": os.path.basename(path), "type": "pptx", "slide_idx": slide_idx, "notes": True}
                        for chunk_idx, sub in enumerate(ResourceIntake.simple_chunker(notes, chunk_words, overlap)):
                            meta = dict(base_meta)
                            meta.update({"chunk_idx": chunk_idx, "excerpt": sub[:200]})
                            results.append({"text": sub, "meta": meta})
            except Exception:
                logger.exception("Failed to read notes for slide %s in %s", slide_idx, path)
        return results

    @staticmethod
    def extract_pdf(path: str, ocr_if_empty: bool = True, chunk_words: int = 200, overlap: int = 40) -> List[Dict[str, Any]]:
        results = []
        try:
            doc = fitz.open(path)
        except Exception:
            logger.exception("Failed to open PDF %s", path)
            return results

        for page_num in range(len(doc)):
            try:
                page = doc[page_num]
                text = page.get_text("text").strip()
                if not text and ocr_if_empty:
                    try:
                        pix = page.get_pixmap(dpi=200)
                        mode = "RGB" if pix.n < 4 else "RGBA"
                        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                        text = pytesseract.image_to_string(img).strip()
                    except Exception:
                        logger.exception("OCR fallback failed for %s page %s", path, page_num + 1)
                        text = ""
                if text:
                    base_meta = {"source": os.path.basename(path), "type": "pdf", "page": page_num + 1}
                    for chunk_idx, sub in enumerate(ResourceIntake.simple_chunker(text, chunk_words, overlap)):
                        meta = dict(base_meta)
                        meta.update({"chunk_idx": chunk_idx, "excerpt": sub[:200]})
                        results.append({"text": sub, "meta": meta})
            except Exception:
                logger.exception("Error extracting page %s from %s", page_num + 1, path)
        return results

    @staticmethod
    def extract_from_path(path: str, **kwargs) -> List[Dict[str, Any]]:
        p = Path(path)
        ext = p.suffix.lower()
        if ext == ".docx":
            return ResourceIntake.extract_docx(path, **kwargs)
        if ext in (".pptx", ".ppt"):
            return ResourceIntake.extract_pptx(path, **kwargs)
        if ext == ".pdf":
            return ResourceIntake.extract_pdf(path, **kwargs)
        logger.warning("Unsupported file type: %s", ext)
        return []
