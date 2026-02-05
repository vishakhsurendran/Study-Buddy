import os
from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract

class ResourceIntake:
    def extract_docx(path: str):
        doc = DocxDocument(path)
        chunks = []
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            meta = {"source": os.path.basename(path), "type":"docx", "para_idx": i}
            chunks.append((text, meta))
        return chunks
    
    def extract_pptx(path: str):
        prs = Presentation(path)
        chunks = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        meta = {"source": os.path.basename(path), "type":"pptx", "slide_idx": slide_idx, "shape_idx": shape_idx}
                        chunks.append((text, meta))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    meta = {"source": os.path.basename(path), "type":"pptx", "slide_idx": slide_idx, "notes": True}
                    chunks.append((notes, meta))
        return chunks
    
    def extract_pdf(path: str, ocr_if_empty=True):
        doc = fitz.open(path)
        chunks = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            if not text and ocr_if_empty:
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img).strip()
            if text: 
                meta = {"source": os.path.basename(path), "type":"pdf", "page": page_num + 1}
                chunks.append((text, meta))
        return chunks
    
    def extract_from_path(path: str):
        ext = Path(path).suffix.lower()
        if ext == ".docx":
            return ResourceIntake.extract_docx(path)
        if ext in (".pptx", ".ppt"):
            return ResourceIntake.extract_pptx(path)
        if ext == ".pdf":
            return ResourceIntake.extract_pdf(path)
        return []
    